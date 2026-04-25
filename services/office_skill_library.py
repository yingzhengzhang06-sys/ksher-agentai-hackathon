"""
Office技能库 — 上传文档→K2.6分析→提取风格模板→复用生成

支持：PPT风格学习、Word写作风格学习
"""
import json
import sqlite3
import time
from dataclasses import dataclass, asdict
from typing import Any

from config import DATA_DIR

SKILL_DB_PATH = f"{DATA_DIR}/skill_library.db"


@dataclass
class OfficeSkill:
    """Office技能模板"""
    skill_id: str
    name: str
    doc_type: str           # "ppt" | "word"
    source_file: str
    created_at: float
    # 风格特征
    style_features: dict    # 配色/字体/版式
    structure_pattern: dict # 章节结构
    content_rules: dict     # 内容组织规则
    example_snippets: list  # 示例片段

    def to_dict(self) -> dict:
        return asdict(self)


class OfficeSkillLibrary:
    """Office技能库管理"""

    def __init__(self):
        self.db_path = SKILL_DB_PATH
        self._init_db()

    def _init_db(self):
        """初始化数据库"""
        with sqlite3.connect(self.db_path) as conn:
            conn.execute("""
                CREATE TABLE IF NOT EXISTS office_skills (
                    skill_id TEXT PRIMARY KEY,
                    name TEXT,
                    doc_type TEXT,
                    source_file TEXT,
                    created_at REAL,
                    style_features TEXT,
                    structure_pattern TEXT,
                    content_rules TEXT,
                    example_snippets TEXT
                )
            """)
            conn.commit()

    def learn_from_document(self, file_bytes: bytes, filename: str,
                           file_type: str, llm_client=None) -> OfficeSkill | None:
        """
        上传文档，K2.6分析提取风格模板。

        Args:
            file_bytes: 文件字节内容
            filename: 文件名
            file_type: "pptx" | "docx"
            llm_client: LLMClient实例（用于K2.6分析）

        Returns:
            OfficeSkill: 提取的技能模板
        """
        # 步骤1：提取文档结构
        if file_type == "pptx":
            structure = self._extract_pptx_structure(file_bytes)
        elif file_type == "docx":
            structure = self._extract_docx_structure(file_bytes)
        else:
            return None

        # 步骤2：K2.6分析（如果提供了LLMClient）
        style_features = {}
        content_rules = {}
        if llm_client:
            try:
                analysis = self._analyze_with_k26(structure, filename, llm_client)
                style_features = analysis.get("style_features", {})
                content_rules = analysis.get("content_rules", {})
            except Exception:
                pass

        # 步骤3：创建技能
        skill = OfficeSkill(
            skill_id=f"sk_{int(time.time())}_{filename[:20]}",
            name=f"{filename}风格",
            doc_type=file_type,
            source_file=filename,
            created_at=time.time(),
            style_features=style_features or self._default_style_features(file_type),
            structure_pattern={"extracted": structure[:2000]},
            content_rules=content_rules or {},
            example_snippets=[structure[:500]],
        )

        # 步骤4：持久化
        self._save_skill(skill)
        return skill

    def _extract_pptx_structure(self, file_bytes: bytes) -> str:
        """提取PPT结构"""
        try:
            from pptx import Presentation
            from io import BytesIO
            prs = Presentation(BytesIO(file_bytes))
            parts = []
            for i, slide in enumerate(prs.slides[:10]):  # 只分析前10页
                texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text.strip()[:100])
                parts.append(f"Slide {i+1}: {' | '.join(texts[:3])}")
            return "\n".join(parts)
        except Exception as e:
            return f"PPT结构提取失败: {e}"

    def _extract_docx_structure(self, file_bytes: bytes) -> str:
        """提取Word结构"""
        try:
            from docx import Document
            from io import BytesIO
            doc = Document(BytesIO(file_bytes))
            parts = []
            for para in doc.paragraphs[:30]:  # 只分析前30段
                if para.text.strip():
                    parts.append(para.text.strip()[:100])
            return "\n".join(parts)
        except Exception as e:
            return f"Word结构提取失败: {e}"

    def _analyze_with_k26(self, structure: str, filename: str,
                          llm_client) -> dict:
        """使用K2.6分析文档风格"""
        system_prompt = """你是Office文档风格分析专家。请分析以下文档结构，提取可复用的风格模板。

输出JSON格式：
{
  "style_features": {
    "color_scheme": "主色调描述",
    "font_style": "字体风格描述",
    "layout_pattern": "版式特点"
  },
  "content_rules": {
    "heading_structure": "标题层级规则",
    "content_density": "内容密度特点",
    "data_presentation": "数据展示方式"
  }
}"""
        user_prompt = f"请分析以下文档的风格特征：\n\n文件名：{filename}\n\n文档结构：\n{structure[:3000]}"

        raw = llm_client.call_sync(
            agent_name="skill_learner",
            system=system_prompt,
            user_msg=user_prompt,
            temperature=0.3,
        )

        try:
            return json.loads(raw)
        except (json.JSONDecodeError, TypeError):
            # 尝试提取JSON
            import re
            m = re.search(r"\{.*\}", raw, re.DOTALL)
            if m:
                try:
                    return json.loads(m.group())
                except:
                    pass
            # 分析失败，返回错误信息让调用方感知
            return {"_analysis_error": "K2.6风格分析失败，使用默认风格", "_raw_preview": raw[:200]}

    def _default_style_features(self, doc_type: str) -> dict:
        """默认风格特征"""
        if doc_type == "pptx":
            return {
                "color_scheme": "Ksher品牌色 (#E83E4C主色 + #00C9A7辅助色)",
                "font_style": "标题粗体32pt，正文16pt",
                "layout_pattern": "封面深色背景 + 内容页白色背景",
            }
        return {
            "color_scheme": "标准商务风格",
            "font_style": "宋体/微软雅黑",
            "layout_pattern": "标准段落排版",
        }

    def _save_skill(self, skill: OfficeSkill):
        """保存技能到数据库"""
        with sqlite3.connect(self.db_path) as conn:
            conn.execute(
                """INSERT OR REPLACE INTO office_skills
                   (skill_id, name, doc_type, source_file, created_at,
                    style_features, structure_pattern, content_rules, example_snippets)
                   VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?)""",
                (
                    skill.skill_id, skill.name, skill.doc_type, skill.source_file,
                    skill.created_at,
                    json.dumps(skill.style_features, ensure_ascii=False),
                    json.dumps(skill.structure_pattern, ensure_ascii=False),
                    json.dumps(skill.content_rules, ensure_ascii=False),
                    json.dumps(skill.example_snippets, ensure_ascii=False),
                )
            )
            conn.commit()

    def list_skills(self, doc_type: str | None = None) -> list[OfficeSkill]:
        """列出所有技能"""
        with sqlite3.connect(self.db_path) as conn:
            if doc_type:
                rows = conn.execute(
                    "SELECT * FROM office_skills WHERE doc_type = ? ORDER BY created_at DESC",
                    (doc_type,)
                ).fetchall()
            else:
                rows = conn.execute(
                    "SELECT * FROM office_skills ORDER BY created_at DESC"
                ).fetchall()

        skills = []
        for row in rows:
            skills.append(OfficeSkill(
                skill_id=row[0],
                name=row[1],
                doc_type=row[2],
                source_file=row[3],
                created_at=row[4],
                style_features=json.loads(row[5]) if row[5] else {},
                structure_pattern=json.loads(row[6]) if row[6] else {},
                content_rules=json.loads(row[7]) if row[7] else {},
                example_snippets=json.loads(row[8]) if row[8] else [],
            ))
        return skills

    def get_skill(self, skill_id: str) -> OfficeSkill | None:
        """获取单个技能"""
        with sqlite3.connect(self.db_path) as conn:
            row = conn.execute(
                "SELECT * FROM office_skills WHERE skill_id = ?", (skill_id,)
            ).fetchone()
        if not row:
            return None
        return OfficeSkill(
            skill_id=row[0],
            name=row[1],
            doc_type=row[2],
            source_file=row[3],
            created_at=row[4],
            style_features=json.loads(row[5]) if row[5] else {},
            structure_pattern=json.loads(row[6]) if row[6] else {},
            content_rules=json.loads(row[7]) if row[7] else {},
            example_snippets=json.loads(row[8]) if row[8] else [],
        )

    def delete_skill(self, skill_id: str) -> bool:
        """删除技能"""
        with sqlite3.connect(self.db_path) as conn:
            cur = conn.execute("DELETE FROM office_skills WHERE skill_id = ?", (skill_id,))
            conn.commit()
            return cur.rowcount > 0


# 单例
_skill_library = None


def get_skill_library() -> OfficeSkillLibrary:
    """获取技能库单例"""
    global _skill_library
    if _skill_library is None:
        _skill_library = OfficeSkillLibrary()
    return _skill_library

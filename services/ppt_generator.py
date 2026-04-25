"""
PPT生成器 — 从结构化大纲生成.pptx文件

使用python-pptx将AI生成的幻灯片大纲渲染为可下载的PowerPoint文件。
支持品牌色注入、自定义布局、演讲者备注。
"""
from io import BytesIO
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from config import BRAND_COLORS


class PPTGenerator:
    """PPT文件生成器"""

    # 品牌色RGB
    PRIMARY = RGBColor(0xE8, 0x3E, 0x4C)      # #E83E4C
    SECONDARY = RGBColor(0x1A, 0x1A, 0x2E)    # #1A1A2E
    ACCENT = RGBColor(0x00, 0xC9, 0xA7)       # #00C9A7
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    DARK_TEXT = RGBColor(0x1D, 0x21, 0x29)
    GRAY_TEXT = RGBColor(0x8A, 0x8F, 0x99)

    def __init__(self):
        self.prs = None

    def generate(self, outline: dict, skill_template: dict | None = None) -> bytes:
        """
        生成.pptx文件字节流。

        Args:
            outline: {
                "title": str,
                "subtitle": str,
                "slides": [
                    {
                        "slide_num": int,
                        "title": str,
                        "content": str,          # 要点列表，\n分隔
                        "speaker_notes": str,
                        "layout": str,           # title, content, two_column, chart
                    }
                ]
            }
            skill_template: 可选，OfficeSkillLibrary提取的风格模板

        Returns:
            bytes: .pptx文件内容
        """
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)

        slides_data = outline.get("slides", [])

        # 封面页
        self._add_title_slide(
            title=outline.get("title", "方案演示"),
            subtitle=outline.get("subtitle", ""),
        )

        # 内容页
        for slide_data in slides_data:
            layout_type = slide_data.get("layout", "content")
            if layout_type == "title":
                self._add_title_slide(
                    title=slide_data.get("title", ""),
                    subtitle=slide_data.get("content", ""),
                )
            elif layout_type == "two_column":
                self._add_two_column_slide(slide_data)
            else:
                self._add_content_slide(slide_data)

        # 结尾页
        self._add_end_slide()

        # 导出为字节
        buffer = BytesIO()
        self.prs.save(buffer)
        buffer.seek(0)
        return buffer.getvalue()

    def _add_title_slide(self, title: str, subtitle: str = ""):
        """添加封面/标题页"""
        blank_layout = self.prs.slide_layouts[6]  # 空白布局
        slide = self.prs.slides.add_slide(blank_layout)

        # 背景色块（左上到右下渐变效果，用矩形模拟）
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0,
            self.prs.slide_width, self.prs.slide_height
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = self.SECONDARY
        bg_shape.line.fill.background()

        # 品牌色装饰条
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, Inches(6.8),
            self.prs.slide_width, Inches(0.7)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = self.PRIMARY
        accent_bar.line.fill.background()

        # 标题
        title_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.5)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = self.WHITE
        p.alignment = PP_ALIGN.LEFT

        # 副标题
        if subtitle:
            sub_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(4.2), Inches(11.7), Inches(1)
            )
            tf = sub_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(20)
            p.font.color.rgb = self.GRAY_TEXT
            p.alignment = PP_ALIGN.LEFT

        # Ksher Logo占位文字
        logo_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(6.9), Inches(3), Inches(0.4)
        )
        tf = logo_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Ksher"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = self.WHITE

    def _add_content_slide(self, slide_data: dict):
        """添加标准内容页"""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)

        # 白色背景
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0,
            self.prs.slide_width, self.prs.slide_height
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.WHITE
        bg.line.fill.background()

        # 顶部装饰条
        top_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0,
            self.prs.slide_width, Inches(0.15)
        )
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = self.PRIMARY
        top_bar.line.fill.background()

        # 标题
        title = slide_data.get("title", "")
        title_box = slide.shapes.add_textbox(
            Inches(0.6), Inches(0.4), Inches(12), Inches(0.8)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self.DARK_TEXT

        # 内容（按行分割为要点）
        content = slide_data.get("content", "")
        content_box = slide.shapes.add_textbox(
            Inches(0.6), Inches(1.4), Inches(12), Inches(5.5)
        )
        tf = content_box.text_frame
        tf.word_wrap = True

        lines = [line.strip() for line in content.split("\n") if line.strip()]
        for i, line in enumerate(lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            # 检测是否为二级标题（以"##"或数字开头）
            if line.startswith("##") or line.startswith("**"):
                clean = line.replace("##", "").replace("**", "").strip()
                p.text = clean
                p.font.size = Pt(20)
                p.font.bold = True
                p.font.color.rgb = self.PRIMARY
                p.space_before = Pt(16)
            elif line.startswith("-") or line.startswith("•"):
                clean = line[1:].strip()
                p.text = f"  • {clean}"
                p.font.size = Pt(16)
                p.font.color.rgb = self.DARK_TEXT
                p.space_before = Pt(6)
            elif line[0:1].isdigit() and "." in line[:3]:
                p.text = f"  {line}"
                p.font.size = Pt(16)
                p.font.color.rgb = self.DARK_TEXT
                p.space_before = Pt(6)
            else:
                p.text = line
                p.font.size = Pt(16)
                p.font.color.rgb = self.DARK_TEXT
                p.space_before = Pt(6)

        # 演讲者备注
        notes = slide_data.get("speaker_notes", "")
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes

    def _add_two_column_slide(self, slide_data: dict):
        """添加双栏内容页"""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)

        # 白色背景
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0,
            self.prs.slide_width, self.prs.slide_height
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.WHITE
        bg.line.fill.background()

        # 顶部装饰条
        top_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0,
            self.prs.slide_width, Inches(0.15)
        )
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = self.PRIMARY
        top_bar.line.fill.background()

        # 标题
        title = slide_data.get("title", "")
        title_box = slide.shapes.add_textbox(
            Inches(0.6), Inches(0.4), Inches(12), Inches(0.8)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self.DARK_TEXT

        # 左栏
        left_content = slide_data.get("left_content", slide_data.get("content", ""))
        left_box = slide.shapes.add_textbox(
            Inches(0.6), Inches(1.4), Inches(5.8), Inches(5.5)
        )
        tf = left_box.text_frame
        tf.word_wrap = True
        lines = [l.strip() for l in left_content.split("\n") if l.strip()]
        for i, line in enumerate(lines[:8]):  # 限制行数
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line.replace("##", "").replace("**", "")
            p.font.size = Pt(15)
            p.font.color.rgb = self.DARK_TEXT
            p.space_before = Pt(6)

        # 右栏
        right_content = slide_data.get("right_content", "")
        if right_content:
            right_box = slide.shapes.add_textbox(
                Inches(6.8), Inches(1.4), Inches(5.8), Inches(5.5)
            )
            tf = right_box.text_frame
            tf.word_wrap = True
            lines = [l.strip() for l in right_content.split("\n") if l.strip()]
            for i, line in enumerate(lines[:8]):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = line.replace("##", "").replace("**", "")
                p.font.size = Pt(15)
                p.font.color.rgb = self.DARK_TEXT
                p.space_before = Pt(6)

    def _add_end_slide(self):
        """添加结尾页（CTA）"""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)

        # 深色背景
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0,
            self.prs.slide_width, self.prs.slide_height
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.SECONDARY
        bg.line.fill.background()

        # CTA标题
        cta_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(2.8), Inches(11.7), Inches(1)
        )
        tf = cta_box.text_frame
        p = tf.paragraphs[0]
        p.text = "选择 Ksher，让跨境收款更简单"
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = self.WHITE
        p.alignment = PP_ALIGN.CENTER

        # 副文案
        sub_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(4.2), Inches(11.7), Inches(0.8)
        )
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = "东南亚本地牌照 | 更低费率 | 更快到账 | 全程合规"
        p.font.size = Pt(18)
        p.font.color.rgb = self.GRAY_TEXT
        p.alignment = PP_ALIGN.CENTER

        # 底部装饰
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, Inches(6.8),
            self.prs.slide_width, Inches(0.7)
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = self.PRIMARY
        accent.line.fill.background()

    def _apply_skill_template(self, slide, template: dict):
        """应用Office技能模板（预留接口）"""
        if not template:
            return
        # TODO: 根据模板调整颜色、字体、布局
        pass
